"""Build small deterministic DOCX, PPTX, and XLSX parser canaries."""

from __future__ import annotations

import argparse
import subprocess
import tempfile
from pathlib import Path


def build_docx(out: Path) -> None:
    html = """<!doctype html><meta charset="utf-8">
<h1>Document parser fixture</h1>
<p>CAPY-DOCX-CANARY-4821 must survive Office conversion and parsing.</p>
<table border="1"><tr><th>Course</th><th>Score</th></tr>
<tr><td>Biology</td><td>92</td></tr><tr><td>History</td><td>88</td></tr></table>
"""
    with tempfile.TemporaryDirectory(prefix="capy-docx-") as temp:
        source = Path(temp) / "office-canary.html"
        source.write_text(html, encoding="utf-8")
        subprocess.run(
            [
                "soffice",
                "--headless",
                "--convert-to",
                "docx:Office Open XML Text",
                "--outdir",
                str(out),
                str(source),
            ],
            check=True,
            capture_output=True,
            text=True,
            timeout=60,
        )


def build_pptx(out: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    deck = Presentation()
    first = deck.slides.add_slide(deck.slide_layouts[1])
    first.shapes.title.text = "Presentation parser fixture"
    first.placeholders[1].text = (
        "CAPY-PPTX-CANARY-5932\n"
        "This slide checks title, body text, and page attribution."
    )
    second = deck.slides.add_slide(deck.slide_layouts[5])
    second.shapes.title.text = "Results table"
    table = second.shapes.add_table(
        3, 2, Inches(1.0), Inches(1.8), Inches(8.0), Inches(2.5)
    ).table
    for row, values in enumerate(
        (
            ("Metric", "Value"),
            ("CAPY-PPTX-TABLE-6043", "present"),
            ("Accuracy", "reviewed"),
        )
    ):
        for column, value in enumerate(values):
            table.cell(row, column).text = value
    deck.save(out / "office-canary.pptx")


def build_xlsx(out: Path) -> None:
    from openpyxl import Workbook
    from openpyxl.styles import Font

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Parser fixture"
    rows = (
        ("CAPY-XLSX-CANARY-7154", "Quarter", "Revenue"),
        ("North", "Q1", 12500),
        ("South", "Q1", 9800),
        ("Total", "", "=SUM(C2:C3)"),
    )
    for row in rows:
        sheet.append(row)
    for cell in sheet[1]:
        cell.font = Font(bold=True)
    sheet.column_dimensions["A"].width = 30
    sheet.column_dimensions["B"].width = 14
    sheet.column_dimensions["C"].width = 18
    sheet.print_area = "A1:C4"
    workbook.save(out / "office-canary.xlsx")


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("out", type=Path)
    args = parser.parse_args()
    args.out.mkdir(parents=True, exist_ok=True)
    build_docx(args.out)
    build_pptx(args.out)
    build_xlsx(args.out)


if __name__ == "__main__":
    main()
